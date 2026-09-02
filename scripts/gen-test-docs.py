#!/usr/bin/env python3
"""生成入库测试用的全类型文档各一份，放入 test-docs/。

覆盖系统允许的 7 种 MIME：PDF / DOCX / XLSX / PPTX / MD / TXT / HTML。
文档共用同一套实体（项目/供应商/人员/制度/部门），便于图谱多跳问答。
"""

from __future__ import annotations

import os
from datetime import date

from docx import Document as WordDoc
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRgb
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt
from reportlab.lib.colors import HexColor, white
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    HRFlowable,
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, "test-docs")

COMPANY = "华辰智造科技股份有限公司"
NAVY = HexColor("#1B365D")
GOLD = HexColor("#B08D57")
INK = HexColor("#1A1A1A")
MUTED = HexColor("#5C6570")
LINE = HexColor("#C9C3B8")
ROW_ALT = HexColor("#F4F1EA")

pdfmetrics.registerFont(TTFont("Hei", "/System/Library/Fonts/STHeiti Light.ttc", subfontIndex=0))
pdfmetrics.registerFont(TTFont("HeiBold", "/System/Library/Fonts/STHeiti Medium.ttc", subfontIndex=0))


def ensure_out():
    os.makedirs(OUT_DIR, exist_ok=True)


def styles():
    return {
        "cover_co": ParagraphStyle("cover_co", fontName="Hei", fontSize=11, textColor=GOLD, alignment=TA_CENTER, leading=16),
        "cover_title": ParagraphStyle("cover_title", fontName="HeiBold", fontSize=20, textColor=NAVY, alignment=TA_CENTER, leading=30, spaceBefore=8, spaceAfter=8),
        "cover_sub": ParagraphStyle("cover_sub", fontName="Hei", fontSize=11, textColor=MUTED, alignment=TA_CENTER, leading=18),
        "h1": ParagraphStyle("h1", fontName="HeiBold", fontSize=14, textColor=NAVY, leading=22, spaceBefore=14, spaceAfter=8),
        "h2": ParagraphStyle("h2", fontName="HeiBold", fontSize=12, textColor=NAVY, leading=18, spaceBefore=10, spaceAfter=6),
        "body": ParagraphStyle("body", fontName="Hei", fontSize=10, textColor=INK, leading=17, alignment=TA_JUSTIFY, spaceAfter=6),
        "meta": ParagraphStyle("meta", fontName="Hei", fontSize=9, textColor=MUTED, leading=14, alignment=TA_LEFT),
        "cell": ParagraphStyle("cell", fontName="Hei", fontSize=8.5, textColor=INK, leading=13, alignment=TA_LEFT),
        "cell_c": ParagraphStyle("cell_c", fontName="Hei", fontSize=8.5, textColor=INK, leading=13, alignment=TA_CENTER),
        "th": ParagraphStyle("th", fontName="HeiBold", fontSize=8.5, textColor=white, leading=13, alignment=TA_CENTER),
        "caption": ParagraphStyle("caption", fontName="Hei", fontSize=8, textColor=MUTED, leading=12, alignment=TA_CENTER, spaceBefore=2, spaceAfter=10),
        "secret": ParagraphStyle("secret", fontName="HeiBold", fontSize=9, textColor=HexColor("#8B1E3F"), alignment=TA_CENTER, leading=14),
    }


S = styles()


def P(text, style="body"):
    return Paragraph(text, S[style])


def bullets(items):
    return ListFlowable(
        [ListItem(P(i), leftIndent=12, bulletColor=NAVY) for i in items],
        bulletType="bullet",
        start="•",
        leftIndent=16,
        bulletFontName="Hei",
        bulletFontSize=9,
        spaceAfter=8,
    )


def table(headers, rows, col_widths=None, caption=None):
    head = [P(h, "th") for h in headers]
    body = [[P(str(c), "cell_c" if i == 0 else "cell") for i, c in enumerate(r)] for r in rows]
    t = Table([head] + body, colWidths=col_widths, repeatRows=1)
    cmds = [
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("GRID", (0, 0), (-1, -1), 0.4, LINE),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]
    for i in range(1, len(body) + 1):
        if i % 2 == 0:
            cmds.append(("BACKGROUND", (0, i), (-1, i), ROW_ALT))
    t.setStyle(TableStyle(cmds))
    flow = [t]
    flow.append(P(caption, "caption") if caption else Spacer(1, 8))
    return flow


def hr():
    return HRFlowable(width="100%", thickness=1.2, color=NAVY, spaceBefore=4, spaceAfter=8)


def header_footer(doc_no, title):
    def _draw(canvas, doc):
        canvas.saveState()
        w, h = A4
        canvas.setFillColor(NAVY)
        canvas.rect(0, h - 16 * mm, w, 16 * mm, fill=1, stroke=0)
        canvas.setFillColor(GOLD)
        canvas.rect(0, h - 16.8 * mm, w, 1.2 * mm, fill=1, stroke=0)
        canvas.setFillColor(white)
        canvas.setFont("Hei", 8)
        canvas.drawString(18 * mm, h - 8.5 * mm, COMPANY)
        canvas.drawRightString(w - 18 * mm, h - 8.5 * mm, "内部资料 · 知识库测试文档")
        canvas.setFillColor(LINE)
        canvas.rect(0, 12 * mm, w, 0.4, fill=1, stroke=0)
        canvas.setFillColor(MUTED)
        canvas.setFont("Hei", 8)
        canvas.drawString(18 * mm, 7 * mm, f"{doc_no}　{title}")
        canvas.drawRightString(w - 18 * mm, 7 * mm, f"第 {doc.page} 页")
        canvas.restoreState()

    return _draw


# ---------------------------------------------------------------------------
# 1. PDF · 差旅费用管理制度
# ---------------------------------------------------------------------------
def gen_pdf():
    path = os.path.join(OUT_DIR, "差旅费用管理制度.pdf")
    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=22 * mm,
        bottomMargin=18 * mm,
        title="差旅费用管理制度",
        author=COMPANY,
    )
    flow = [
        P("机密等级：内部公开", "secret"),
        Spacer(1, 6),
        P(COMPANY, "cover_co"),
        P("财务部 · 制度文件", "cover_co"),
        hr(),
        P("差旅费用管理制度", "cover_title"),
        hr(),
        P("与《供应商付款审批制度》《项目预算与成本核算办法》配套执行", "cover_sub"),
        Spacer(1, 8),
        P("文档编号：HC-FIN-2025-012　发布部门：财务部　生效日期：2025-03-01　批准人：林婉清（财务总监）", "meta"),
        Spacer(1, 6),
        P("一、目的与适用范围", "h1"),
        P(
            "为规范华辰智造科技股份有限公司员工因公出差的费用标准与报销流程，控制差旅成本，"
            "特制定本制度。本制度适用于全体正式员工、劳务派遣人员，以及星云ERP升级项目、"
            "智能工厂一期项目、海外供应链数字化项目等在建项目的驻场人员。",
        ),
        P("二、管理职责", "h1"),
        P("财务部（负责人林婉清）负责标准制定、报销审核与账务处理。各部门负责人对本部门出差必要性负责。项目管理办公室（负责人周嘉宁）对项目差旅预算占用进行复核。审计部对异常报销进行抽查。"),
        P("三、城市分级与住宿标准", "h1"),
        P("住宿费按城市分级据实报销，超额自理。同一城市连续住宿超过 14 天，自第 15 日起按标准的 80% 执行。"),
        *table(
            ["城市分级", "覆盖城市（示例）", "总监及以上", "经理级", "普通员工"],
            [
                ["一类", "北京、上海、深圳、广州", "900 元/晚", "700 元/晚", "550 元/晚"],
                ["二类", "杭州、南京、成都、武汉、苏州", "700 元/晚", "550 元/晚", "420 元/晚"],
                ["三类", "其他省会及计划单列市", "550 元/晚", "420 元/晚", "320 元/晚"],
                ["四类", "地级市及以下", "400 元/晚", "320 元/晚", "260 元/晚"],
            ],
            col_widths=[28 * mm, 52 * mm, 30 * mm, 30 * mm, 30 * mm],
            caption="表 1　住宿费限额（含服务费，不含早餐；单位：人民币）",
        ),
        P("四、交通与餐饮", "h1"),
        P("2. 交通：总监及以上可乘坐飞机公务舱或高铁一等座；经理级及以下乘坐经济舱或二等座。单程票价超过 2,000 元须事先经财务部林婉清审批。市内交通优先网约车企业账户，日限额 150 元。"),
        P("3. 餐饮：一类城市餐补 150 元/人·天，二类 120 元，三类及以下 100 元。已由接待方安排工作餐的，当日餐补减半。不得另外报销酒水。"),
        P("五、项目差旅特别规定", "h1"),
        P(
            "星云ERP升级项目（项目负责人周嘉宁，主要供应商华云科技、天枢软件）驻场苏州期间，"
            "住宿统一由项目管理办公室签订协议酒店，不再按表 1 个人限额报销。"
            "智能工厂一期项目（项目负责人吴振华，咨询供应商锐智咨询）赴东莞产线出差，"
            "交通由远航物流提供班车的，不得重复报销市内交通。"
            "海外供应链数字化项目（项目负责人孙雨桐）出境差旅另行适用《因公出国（境）管理办法》，"
            "本制度仅覆盖境内段。",
        ),
        P("六、报销流程", "h1"),
        bullets(
            [
                "出差前在 OA 提交差旅申请，部门负责人审批；预计费用超过 8,000 元须财务部会签。",
                "返回后 10 个工作日内在费控系统提交报销单，附行程单、发票、会议通知或项目驻场证明。",
                "金额 5,000 元以下由部门负责人终审；5,000–20,000 元由财务部林婉清终审；超过 20,000 元由总经理陈思远审批。",
                "项目差旅须同时占用对应项目预算，由项目管理办公室周嘉宁或其授权人确认预算科目。",
                "违反本制度的超标费用一律自理；虚假报销移交审计部并按《员工奖惩办法》处理。",
            ]
        ),
        P("七、附则", "h1"),
        P("本制度由财务部解释，自 2025 年 3 月 1 日起施行，原《差旅报销管理办法（2023 版）》同时废止。与《供应商付款审批制度》冲突时，对公付款以付款制度为准，个人报销以本制度为准。"),
    ]
    doc.build(flow, onFirstPage=header_footer("HC-FIN-2025-012", "差旅费用管理制度"), onLaterPages=header_footer("HC-FIN-2025-012", "差旅费用管理制度"))
    return path


# ---------------------------------------------------------------------------
# 2. DOCX · 供应商准入与付款管理办法
# ---------------------------------------------------------------------------
def _set_run_font(run, size=11, bold=False, color=None, name="微软雅黑"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if color:
        run.font.color.rgb = RGBColor(*color)


def _add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.color.rgb = RGBColor(0x1B, 0x36, 0x5D)
        run.font.name = "微软雅黑"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    return p


def gen_docx():
    path = os.path.join(OUT_DIR, "供应商准入与付款管理办法.docx")
    doc = WordDoc()
    section = doc.sections[0]
    section.left_margin = Cm(2.2)
    section.right_margin = Cm(2.2)

    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cap.add_run(COMPANY + "　采购部 / 财务部联合发文")
    _set_run_font(r, 12, color=(0xB0, 0x8D, 0x57))

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("供应商准入与付款管理办法")
    _set_run_font(r, 22, bold=True, color=(0x1B, 0x36, 0x5D))

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = meta.add_run("文档编号 HC-PUR-2025-008　生效日期 2025-04-01　批准人：赵启明（采购总监）、林婉清（财务总监）")
    _set_run_font(r, 9, color=(0x5C, 0x65, 0x70))

    _add_heading(doc, "1  目的", 1)
    p = doc.add_paragraph(
        "规范华辰智造科技股份有限公司合格供应商的准入、分级、付款与退出，"
        "降低供应中断与资金风险。本制度与《差旅费用管理制度》《项目预算与成本核算办法》配套执行。"
    )
    for run in p.runs:
        _set_run_font(run, 11)

    _add_heading(doc, "2  职责分工", 1)
    p = doc.add_paragraph(
        "采购部（负责人赵启明）负责供应商寻源、准入评审与合同签订。"
        "财务部（负责人林婉清）负责付款审核、账期管理与对账。"
        "项目管理办公室（负责人周嘉宁）对项目类采购的预算占用与验收签字负责。"
        "审计部对关联交易与超账期付款进行专项检查。"
        "总经理陈思远对年度框架协议及单笔超过 200 万元的合同拥有终审权。"
    )
    for run in p.runs:
        _set_run_font(run, 11)

    _add_heading(doc, "3  准入条件", 1)
    for item in [
        "依法设立满三年，近一年无重大违法记录，注册资本不低于人民币 500 万元。",
        "具备与供应品类匹配的资质：软件类须等保/ISO 27001；物流类须道路运输经营许可；咨询类须相关行业认证。",
        "须通过采购部组织的现场或远程尽调，评审小组至少包含采购部、财务部、需求部门各一人。",
        "关联方供应商（含持股 5% 以上）须额外提交《关联交易披露表》，报审计部备案后方可入库。",
    ]:
        p = doc.add_paragraph(item, style="List Number")
        for run in p.runs:
            _set_run_font(run, 11)

    _add_heading(doc, "4  现有合格供应商与服务项目", 1)
    p = doc.add_paragraph("截至 2025 年 6 月，公司合格供应商名录中与在建项目直接相关的核心供应商如下。")
    for run in p.runs:
        _set_run_font(run, 11)

    table = doc.add_table(rows=5, cols=5)
    table.style = "Table Grid"
    headers = ["供应商", "主联系人", "服务项目", "项目负责人", "账期"]
    rows = [
        ["华云科技", "钱志远（客户成功）", "星云ERP升级项目、海外供应链数字化项目", "周嘉宁 / 孙雨桐", "月结 45 天"],
        ["天枢软件", "何嘉慧", "星云ERP升级项目（MES 对接模块）", "周嘉宁", "验收后 30 天"],
        ["远航物流", "郑浩然（商务）", "智能工厂一期项目、海外供应链数字化项目", "吴振华 / 孙雨桐", "月结 30 天"],
        ["锐智咨询", "黄晓薇（项目经理）", "智能工厂一期项目（精益咨询）", "吴振华", "里程碑付款"],
    ]
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for run in cell.paragraphs[0].runs:
            _set_run_font(run, 9, bold=True, color=(255, 255, 255))
        shading = cell._tc.get_or_add_tcPr()
        shading.append(parse_xml(r'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" w:fill="1B365D"/>'))
    for ridx, row in enumerate(rows, start=1):
        for cidx, val in enumerate(row):
            cell = table.rows[ridx].cells[cidx]
            cell.text = val
            for run in cell.paragraphs[0].runs:
                _set_run_font(run, 9)

    _add_heading(doc, "5  付款审批权限", 1)
    p = doc.add_paragraph(
        "单笔 10 万元以下：需求部门负责人审批，财务部复核后支付。"
        "单笔 10–50 万元：采购部赵启明与财务部林婉清会签。"
        "单笔 50–200 万元：增加项目管理办公室周嘉宁会签（项目类）或总经理陈思远审批（非项目类）。"
        "超过 200 万元或涉及华云科技、天枢软件等战略供应商的框架追加，必须经总经理陈思远书面批准。"
        "预付款比例原则上不超过合同额的 30%；锐智咨询按里程碑付款，不得预付。"
    )
    for run in p.runs:
        _set_run_font(run, 11)

    _add_heading(doc, "6  禁止事项", 1)
    for item in [
        "未入库供应商不得签署合同，紧急采购须 5 个工作日内补办准入。",
        "不得将星云ERP升级项目款项挪用于智能工厂一期项目或其他项目。",
        "个人垫付的对公费用不得走《差旅费用管理制度》渠道报销，须走对公付款。",
    ]:
        p = doc.add_paragraph(item, style="List Bullet")
        for run in p.runs:
            _set_run_font(run, 11)

    _add_heading(doc, "7  附则", 1)
    p = doc.add_paragraph("本制度由采购部解释，自 2025 年 4 月 1 日起施行。供应商分级评分细则见附件《2025年度合格供应商名录》。")
    for run in p.runs:
        _set_run_font(run, 11)

    doc.save(path)
    return path


# ---------------------------------------------------------------------------
# 3. XLSX · 2025年度合格供应商名录
# ---------------------------------------------------------------------------
def gen_xlsx():
    path = os.path.join(OUT_DIR, "2025年度合格供应商名录.xlsx")
    wb = Workbook()
    header_font = Font(name="微软雅黑", bold=True, color="FFFFFF", size=10)
    cell_font = Font(name="微软雅黑", size=10)
    header_fill = PatternFill("solid", fgColor="1B365D")
    alt_fill = PatternFill("solid", fgColor="F4F1EA")
    thin = Border(
        left=Side(style="thin", color="C9C3B8"),
        right=Side(style="thin", color="C9C3B8"),
        top=Side(style="thin", color="C9C3B8"),
        bottom=Side(style="thin", color="C9C3B8"),
    )
    wrap = Alignment(wrap_text=True, vertical="center")

    def style_sheet(ws, headers, rows, widths):
        ws.append(headers)
        for row in rows:
            ws.append(row)
        for col, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(col)].width = w
        for cell in ws[1]:
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin
        for r in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=len(headers)):
            for i, cell in enumerate(r):
                cell.font = cell_font
                cell.alignment = wrap
                cell.border = thin
                if cell.row % 2 == 0:
                    cell.fill = alt_fill
        ws.row_dimensions[1].height = 22
        ws.auto_filter.ref = ws.dimensions
        ws.freeze_panes = "A2"
        ws.sheet_properties.pageSetUpPr.fitToPage = True

    ws1 = wb.active
    ws1.title = "合格供应商清单"
    style_sheet(
        ws1,
        ["供应商编码", "供应商名称", "品类", "主联系人", "联系部门", "准入日期", "等级", "账期", "年度额度(万元)", "状态"],
        [
            ["V-2023-018", "华云科技", "企业软件 / ERP 实施", "钱志远", "采购部", "2023-05-12", "A 战略", "月结 45 天", 860, "在用"],
            ["V-2024-006", "天枢软件", "MES / 工业软件", "何嘉慧", "采购部", "2024-02-20", "B 优选", "验收后 30 天", 320, "在用"],
            ["V-2022-041", "远航物流", "仓储物流 / 国际货代", "郑浩然", "采购部", "2022-11-08", "A 战略", "月结 30 天", 540, "在用"],
            ["V-2024-019", "锐智咨询", "精益与智能制造咨询", "黄晓薇", "采购部", "2024-07-01", "B 优选", "里程碑付款", 180, "在用"],
            ["V-2021-003", "金桥融资租赁", "设备融资租赁", "孙雨桐（对接）", "财务部", "2021-03-15", "C 合格", "季度结算", 1200, "在用"],
            ["V-2020-011", "鼎信会计师事务所", "审计与内控咨询", "钱志远", "审计部", "2020-09-01", "B 优选", "项目结束 15 天", 90, "在用"],
        ],
        [14, 18, 22, 14, 12, 14, 12, 16, 16, 10],
    )

    ws2 = wb.create_sheet("项目服务对照")
    style_sheet(
        ws2,
        ["项目名称", "项目负责人", "所属部门", "供应商", "服务内容", "合同额(万元)", "付款节点", "适用制度"],
        [
            ["星云ERP升级项目", "周嘉宁", "项目管理办公室", "华云科技", "ERP 核心模块实施与数据迁移", 480, "上线 40% + 验收 60%", "供应商付款审批制度"],
            ["星云ERP升级项目", "周嘉宁", "项目管理办公室", "天枢软件", "MES 对接与车间报工接口", 160, "验收后 30 天全额", "供应商付款审批制度"],
            ["智能工厂一期项目", "吴振华", "研发中心", "锐智咨询", "精益产线规划与数字化蓝图", 95, "蓝图/试运行/验收三期", "项目预算与成本核算办法"],
            ["智能工厂一期项目", "吴振华", "研发中心", "远航物流", "产线设备国内干线运输与仓储", 70, "月结 30 天", "供应商付款审批制度"],
            ["海外供应链数字化项目", "孙雨桐", "采购部", "华云科技", "供应商协同门户与海外仓系统", 210, "月结 45 天", "供应商付款审批制度"],
            ["海外供应链数字化项目", "孙雨桐", "采购部", "远航物流", "东南亚专线与海外仓操作", 155, "月结 30 天", "供应商付款审批制度"],
            ["2025年度全面预算编制项目", "林婉清", "财务部", "锐智咨询", "预算模型辅导（兼职）", 28, "里程碑付款", "项目预算与成本核算办法"],
            ["应收账款专项清收项目", "林婉清", "财务部", "鼎信会计师事务所", "账龄复核与内控评价", 45, "报告提交后 15 天", "内部控制评价办法"],
        ],
        [24, 12, 16, 16, 32, 14, 22, 24],
    )

    ws3 = wb.create_sheet("编制说明")
    ws3["A1"] = f"{COMPANY}　2025 年度合格供应商名录"
    ws3["A1"].font = Font(name="微软雅黑", bold=True, size=14, color="1B365D")
    notes = [
        "编制部门：采购部（负责人赵启明）；会签：财务部林婉清、项目管理办公室周嘉宁。",
        "批准人：总经理陈思远。发布日期：2025-06-30。密级：内部公开。",
        "本表是《供应商准入与付款管理办法》的附件，入库时请保持供应商、项目、人员名称与制度原文一致。",
        "华云科技同时服务星云ERP升级项目与海外供应链数字化项目，是图谱多跳的关键供应商。",
        "远航物流同时服务智能工厂一期项目与海外供应链数字化项目。",
        "差旅类个人费用不在本表核算，适用《差旅费用管理制度》。",
        "等级说明：A 战略（可签框架、可预付）、B 优选（常规采购）、C 合格（一事一议）。",
    ]
    for i, text in enumerate(notes, start=3):
        ws3[f"A{i}"] = text
        ws3[f"A{i}"].font = Font(name="微软雅黑", size=11)
        ws3[f"A{i}"].alignment = Alignment(wrap_text=True)
        ws3.row_dimensions[i].height = 22
    ws3.column_dimensions["A"].width = 92

    wb.save(path)
    return path


# ---------------------------------------------------------------------------
# 4. PPTX · 星云ERP升级项目中期汇报
# ---------------------------------------------------------------------------
NAVY_RGB = PptRgb(0x1B, 0x36, 0x5D)
GOLD_RGB = PptRgb(0xB0, 0x8D, 0x57)
WHITE = PptRgb(0xFF, 0xFF, 0xFF)
INK_RGB = PptRgb(0x1A, 0x1A, 0x1A)


def _add_bg(slide, rgb=NAVY_RGB):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _textbox(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return box


def _bullets(slide, l, t, w, h, items, size=16, color=INK_RGB):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = PptPt(8)
    return box


def gen_pptx():
    path = os.path.join(OUT_DIR, "星云ERP升级项目中期汇报.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 封面
    s = prs.slides.add_slide(blank)
    _add_bg(s, NAVY_RGB)
    _textbox(s, 0.8, 1.6, 11, 0.4, COMPANY + "  ·  项目管理办公室", 16, False, GOLD_RGB)
    _textbox(s, 0.8, 2.2, 11.5, 1.2, "星云ERP升级项目中期汇报", 36, True, WHITE)
    _textbox(s, 0.8, 3.5, 11, 0.5, "2025 年半年度经营会　｜　内部资料", 18, False, GOLD_RGB)
    _textbox(
        s,
        0.8,
        5.6,
        11,
        1.0,
        "项目负责人：周嘉宁（项目管理办公室）\n汇报日期：2025-06-18　　批准：总经理陈思远",
        14,
        False,
        WHITE,
    )

    # 项目概览
    s = prs.slides.add_slide(blank)
    _textbox(s, 0.6, 0.3, 12, 0.6, "01  项目概览", 26, True, NAVY_RGB)
    _bullets(
        s,
        0.7,
        1.1,
        12,
        5.8,
        [
            "项目名称：星云ERP升级项目。目标是替换 2018 年版用友系统，打通财务、采购、生产、库存。",
            "项目负责人：周嘉宁（项目管理办公室）。技术负责人：研发中心郑浩然。业务负责人：财务部林婉清。",
            "核心供应商：华云科技（ERP 实施，合同额 480 万元，账期月结 45 天）；天枢软件（MES 对接，160 万元）。",
            "预算总额 720 万元，受《项目预算与成本核算办法》约束；对公付款执行《供应商付款审批制度》。",
            "周期：2025-01-06 启动，计划 2025-11-30 上线。当前进度 58%，落后基线约 2 周。",
            "协同项目：海外供应链数字化项目（负责人孙雨桐）将复用华云科技的供应商协同门户模块。",
        ],
        18,
    )

    # 供应商与分工
    s = prs.slides.add_slide(blank)
    _textbox(s, 0.6, 0.3, 12, 0.6, "02  供应商分工与依赖", 26, True, NAVY_RGB)
    rows = [
        ["角色", "单位 / 人员", "工作包", "状态"],
        ["实施总包", "华云科技 / 钱志远", "财务、采购、库存模块；数据迁移", "开发中"],
        ["MES 对接", "天枢软件 / 何嘉慧", "车间报工接口、与智能工厂一期对齐", "联调中"],
        ["物流主数据", "远航物流（数据提供）", "承运商主数据、运价接口样例", "已提供"],
        ["业主团队", "周嘉宁、林婉清、赵启明、郑浩然", "需求确认、UAT、权限矩阵", "进行中"],
    ]
    table = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.2), Inches(12.1), Inches(3.6)).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.4)
    table.columns[2].width = Inches(4.6)
    table.columns[3].width = Inches(2.3)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(14)
                p.font.name = "微软雅黑"
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else INK_RGB
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_RGB if r == 0 else (PptRgb(0xF4, 0xF1, 0xEA) if r % 2 == 0 else WHITE)

    _textbox(s, 0.6, 5.2, 12, 1.6, "风险：华云科技同期还在交付海外供应链数字化项目，资源冲突需采购部赵启明协调；若延期将影响孙雨桐侧门户上线。天枢软件接口依赖智能工厂一期项目（吴振华）的产线编码规则，两边须每周对齐。", 16, False, INK_RGB)

    # 预算与制度
    s = prs.slides.add_slide(blank)
    _textbox(s, 0.6, 0.3, 12, 0.6, "03  预算执行与合规", 26, True, NAVY_RGB)
    _bullets(
        s,
        0.7,
        1.1,
        12,
        5.8,
        [
            "已付款：华云科技 192 万元（合同 40% 开工款，林婉清已签，陈思远备案）。",
            "待付款：天枢软件首期 64 万元，待 UAT 通过后按《供应商付款审批制度》由赵启明、林婉清会签。",
            "差旅：实施团队驻场苏州协议酒店，适用《差旅费用管理制度》项目差旅特别规定，不按个人限额。",
            "禁止：不得将本项目预算拆借给智能工厂一期项目；关联交易披露按《关联交易披露指引》执行。",
            "审计关注点：华云科技同时服务两个项目，需分项目开票，审计部将在三季度抽查。",
        ],
        18,
    )

    # 下一步
    s = prs.slides.add_slide(blank)
    _textbox(s, 0.6, 0.3, 12, 0.6, "04  下一步与决策事项", 26, True, NAVY_RGB)
    _bullets(
        s,
        0.7,
        1.1,
        12,
        5.8,
        [
            "请总经理陈思远批准将上线窗口从 11 月 30 日调整为 12 月 15 日，并追加应急预算 40 万元。",
            "请采购部赵启明与华云科技确认海外供应链数字化项目的资源隔离承诺函。",
            "请研发中心吴振华在 7 月 5 日前冻结智能工厂一期项目的产线编码，供天枢软件使用。",
            "请财务部林婉清预留 12 月验收款，避免与年度全面预算编制项目抢占付款窗口。",
            "下次汇报：2025-08-20 项目管理办公室周例会。",
        ],
        18,
    )

    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# 5. Markdown · 智能工厂一期项目实施 SOP
# ---------------------------------------------------------------------------
def gen_md():
    path = os.path.join(OUT_DIR, "智能工厂一期项目实施SOP.md")
    content = f"""# 智能工厂一期项目实施 SOP

> {COMPANY}　研发中心　文档编号 HC-RD-2025-021　版本 V1.2　生效日期 2025-05-08
> 项目负责人：吴振华　会签：项目管理办公室周嘉宁、采购部赵启明、财务部林婉清

## 1. 目的

本 SOP 规定智能工厂一期项目从蓝图、设备进场、产线联调到验收的操作标准，确保与星云ERP升级项目的主数据一致，并受《项目预算与成本核算办法》《供应商付款审批制度》约束。

## 2. 组织与职责

| 角色 | 姓名 | 部门 | 职责 |
| --- | --- | --- | --- |
| 项目负责人 | 吴振华 | 研发中心 | 范围、进度、技术决策 |
| PMO 监理 | 周嘉宁 | 项目管理办公室 | 预算占用、跨项目协同 |
| 采购接口 | 赵启明 | 采购部 | 供应商合同与到货 |
| 财务接口 | 林婉清 | 财务部 | 付款与成本归集 |
| 精益咨询 | 黄晓薇（供应商侧） | 锐智咨询 | 产线布局与节拍设计 |
| 物流执行 | 郑浩然（供应商侧） | 远航物流 | 设备运输与仓储 |

## 3. 供应商工作包

1. **锐智咨询**：输出精益产线规划、数字化蓝图、试运行评估报告。合同额 95 万元，按蓝图 / 试运行 / 验收三期付款，**不得预付**。
2. **远航物流**：负责东莞基地设备国内干线运输、暂存与开箱。合同额 70 万元，月结 30 天。若使用远航物流班车，驻场人员不得再按《差旅费用管理制度》报销市内交通。
3. **天枢软件**（星云ERP升级项目分包）：提供车间报工接口。编码规则必须以本项目冻结的产线编码为准，由吴振华于每周三与周嘉宁、郑浩然对齐。

华云科技不直接实施本项目，但海外供应链数字化项目（孙雨桐）与本项目共享远航物流主数据，变更须同步。

## 4. 实施阶段

### 4.1 蓝图阶段（已完成）

- 输出物：《东莞一厂精益布局图》《设备清单 V3》。
- 锐智咨询黄晓薇组织工作坊，研发中心、生产、采购部赵启明参加。
- 蓝图需总经理陈思远签字后方可进入采购。

### 4.2 设备进场

- 到货验收单由吴振华与远航物流郑浩然共同签署，扫描件 24 小时内交财务部林婉清。
- 损坏或短少超过 5,000 元，按《供应商付款审批制度》暂扣当期运费。

### 4.3 联调与 ERP 对齐

- 产线编码一旦冻结，星云ERP升级项目不得单方面修改物料与工作中心主数据。
- 异常升级须提交变更单，周嘉宁、吴振华双签。

### 4.4 验收

- 验收委员会：吴振华、周嘉宁、林婉清、赵启明，必要时邀请审计部。
- 锐智咨询尾款在验收报告归档后 10 个工作日内支付。

## 5. 安全与合规

- 驻场必须遵守东莞基地 EHS 规定；加班餐补按《差旅费用管理制度》三类城市标准执行。
- 本项目预算不得拆借给星云ERP升级项目或海外供应链数字化项目。
- 关联交易（如有）按《关联交易披露指引》报审计部备案。

## 6. 参考文件

- 《差旅费用管理制度》（HC-FIN-2025-012）
- 《供应商准入与付款管理办法》（HC-PUR-2025-008）
- 《2025年度合格供应商名录》
- 《星云ERP升级项目中期汇报》
"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    return path


# ---------------------------------------------------------------------------
# 6. TXT · 海外供应链数字化项目周例会纪要
# ---------------------------------------------------------------------------
def gen_txt():
    path = os.path.join(OUT_DIR, "海外供应链数字化项目周例会纪要.txt")
    content = f"""{COMPANY}
海外供应链数字化项目　周例会纪要
编号：HC-PMO-MOM-2025-0624
会议时间：2025年6月24日 14:00-15:30
地点：总部 8 楼第一会议室 / 腾讯会议 638-219-445
主持人：孙雨桐（项目负责人，采购部）
记录人：项目管理办公室 周嘉宁
密级：内部公开

一、出席人员
孙雨桐（采购部 / 项目负责人）
周嘉宁（项目管理办公室）
赵启明（采购部总监）
林婉清（财务部总监）
吴振华（研发中心 / 智能工厂一期项目负责人，列席）
钱志远（华云科技 客户成功经理）
郑浩然（远航物流 商务经理）
陈思远（总经理，后半场听取决策事项）

请假：黄晓薇（锐智咨询，本项目无当期任务）

二、项目进展
1. 华云科技负责的供应商协同门户已完成需求冻结，开发进度 62%。钱志远确认：门户将复用星云ERP升级项目的供应商主数据，避免两套编码。
2. 远航物流东南亚专线与胡志明海外仓 WMS 接口联调延期 1 周，原因是智能工厂一期项目同期占用了远航物流的实施工程师。吴振华说明东莞设备到货高峰将在 7 月 10 日前结束，之后资源释放。
3. 财务侧：林婉清提醒华云科技须按项目分别开票，不得将星云ERP升级项目与本项目费用合并。已付款 84 万元，剩余按月结 45 天执行《供应商付款审批制度》。

三、议题与决议
议题 1：华云科技资源冲突
决议：采购部赵启明于 6 月 27 日前取得华云科技书面承诺，保证星云ERP升级项目与本项目各配备独立实施经理。未完成前，单笔追加采购暂停。

议题 2：与智能工厂一期的物流主数据
决议：以吴振华冻结的承运商编码为唯一标准，远航物流郑浩然在 7 月 3 日前完成映射表。周嘉宁纳入跨项目依赖看板。

议题 3：出境差旅
孙雨桐、钱志远拟于 7 月中旬赴胡志明仓验收。境内段适用《差旅费用管理制度》；境外段走因公出国流程，预算 6.8 万元，需陈思远批准。

议题 4：预算预警
本项目预算 380 万元，已承诺 365 万元。周嘉宁建议冻结非关键需求。林婉清不同意再占用 2025 年度全面预算编制项目的咨询额度。

四、待办
1. 赵启明：华云科技资源隔离承诺函。截止 6-27。
2. 郑浩然 / 吴振华：承运商编码映射表。截止 7-03。
3. 孙雨桐：出境申请呈陈思远。截止 6-28。
4. 林婉清：核对华云科技两项目开票清单，抄送审计部。截止 7-05。
5. 周嘉宁：更新跨项目风险登记册，同步星云ERP升级项目中期汇报材料。截止 6-26。

五、下次会议
2025年7月1日 14:00，同一会议室。请天枢软件何嘉慧列席（MES 报工是否对海外仓出库回传）。

—— 纪要结束，请出席人 24 小时内回复确认 ——
"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    return path


# ---------------------------------------------------------------------------
# 7. HTML · 组织架构与部门负责人一览
# ---------------------------------------------------------------------------
def gen_html():
    path = os.path.join(OUT_DIR, "组织架构与部门负责人一览.html")
    content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <title>组织架构与部门负责人一览 — {COMPANY}</title>
  <style>
    body {{ font-family: "PingFang SC", "Microsoft YaHei", sans-serif; margin: 0; color: #1a1a1a; background: #f7f5f0; }}
    header {{ background: #1B365D; color: #fff; padding: 28px 40px 24px; }}
    header .co {{ color: #B08D57; font-size: 13px; letter-spacing: 1px; }}
    h1 {{ margin: 8px 0 4px; font-size: 26px; }}
    .meta {{ color: #c9c3b8; font-size: 13px; }}
    main {{ max-width: 960px; margin: 24px auto 48px; background: #fff; padding: 32px 40px; box-shadow: 0 1px 4px rgba(0,0,0,.06); }}
    h2 {{ color: #1B365D; border-left: 4px solid #B08D57; padding-left: 10px; }}
    table {{ width: 100%; border-collapse: collapse; margin: 12px 0 24px; font-size: 14px; }}
    th {{ background: #1B365D; color: #fff; padding: 8px 10px; text-align: left; }}
    td {{ border: 1px solid #C9C3B8; padding: 8px 10px; }}
    tr:nth-child(even) td {{ background: #F4F1EA; }}
    .note {{ color: #5C6570; font-size: 13px; line-height: 1.7; }}
    footer {{ text-align: center; color: #5C6570; font-size: 12px; padding: 16px; }}
  </style>
</head>
<body>
  <header>
    <div class="co">{COMPANY}　人力资源部</div>
    <h1>组织架构与部门负责人一览</h1>
    <div class="meta">文档编号 HC-HR-2025-003　发布日期 2025-03-15　编制：黄晓薇（人力资源部）　批准：陈思远（总经理）</div>
  </header>
  <main>
    <h2>1. 高管与部门负责人</h2>
    <p>本表是知识库人员实体的规范名来源。问答与图谱抽取须使用下表中文姓名与部门全称，避免「财务 / Finance / 财会」等别称并存。</p>
    <table>
      <thead>
        <tr><th>职务</th><th>姓名</th><th>部门</th><th>分管范围</th><th>相关制度 / 项目</th></tr>
      </thead>
      <tbody>
        <tr><td>总经理</td><td>陈思远</td><td>总经理办公室</td><td>经营决策、重大合同终审</td><td>单笔超 200 万元付款；出境差旅</td></tr>
        <tr><td>财务总监</td><td>林婉清</td><td>财务部</td><td>核算、资金、差旅与付款审核</td><td>《差旅费用管理制度》；全面预算编制项目</td></tr>
        <tr><td>采购总监</td><td>赵启明</td><td>采购部</td><td>供应商准入、合同、采购执行</td><td>《供应商准入与付款管理办法》</td></tr>
        <tr><td>PMO 主任</td><td>周嘉宁</td><td>项目管理办公室</td><td>项目组合、预算占用、跨项目协同</td><td>星云ERP升级项目负责人</td></tr>
        <tr><td>研发总监</td><td>吴振华</td><td>研发中心</td><td>产品研发、智能制造技术</td><td>智能工厂一期项目负责人</td></tr>
        <tr><td>供应链经理</td><td>孙雨桐</td><td>采购部</td><td>海外仓与供应商协同</td><td>海外供应链数字化项目负责人</td></tr>
        <tr><td>人力资源总监</td><td>黄晓薇</td><td>人力资源部</td><td>编制、招聘、组织发展</td><td>本文件；兼对接锐智咨询驻场</td></tr>
        <tr><td>审计经理</td><td>钱志远</td><td>审计部</td><td>内控评价、关联交易检查</td><td>《内部控制评价办法》</td></tr>
      </tbody>
    </table>
    <p class="note">说明：供应商侧同名人员（华云科技客户成功经理钱志远、锐智咨询项目经理黄晓薇、远航物流商务郑浩然）不属于公司编制，图谱中类型为 Person，但关系应指向 Supplier 而非 Department。</p>

    <h2>2. 在建项目与负责人对照</h2>
    <table>
      <thead>
        <tr><th>项目</th><th>负责人</th><th>主管部门</th><th>核心供应商</th></tr>
      </thead>
      <tbody>
        <tr><td>星云ERP升级项目</td><td>周嘉宁</td><td>项目管理办公室</td><td>华云科技、天枢软件</td></tr>
        <tr><td>智能工厂一期项目</td><td>吴振华</td><td>研发中心</td><td>锐智咨询、远航物流</td></tr>
        <tr><td>海外供应链数字化项目</td><td>孙雨桐</td><td>采购部</td><td>华云科技、远航物流</td></tr>
        <tr><td>2025年度全面预算编制项目</td><td>林婉清</td><td>财务部</td><td>锐智咨询（兼职辅导）</td></tr>
        <tr><td>应收账款专项清收项目</td><td>林婉清</td><td>财务部</td><td>鼎信会计师事务所</td></tr>
      </tbody>
    </table>

    <h2>3. 发文部门与制度归属</h2>
    <ul>
      <li>财务部林婉清：《差旅费用管理制度》《项目预算与成本核算办法》</li>
      <li>采购部赵启明：《供应商准入与付款管理办法》及附件《2025年度合格供应商名录》</li>
      <li>研发中心吴振华：《智能工厂一期项目实施 SOP》</li>
      <li>项目管理办公室周嘉宁：项目汇报与会议纪要类过程文档</li>
      <li>人力资源部黄晓薇：本组织架构文件，每季度更新一次</li>
    </ul>
    <p>联系人力资源部邮箱 hr@huachen-example.local。组织调整须在生效后 5 个工作日内更新本页，否则知识库问答可能引用过期负责人。</p>
  </main>
  <footer>{COMPANY}　内部公开　请勿外传　{date.today().isoformat()}</footer>
</body>
</html>
"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    return path


def main():
    ensure_out()
    paths = [gen_pdf(), gen_docx(), gen_xlsx(), gen_pptx(), gen_md(), gen_txt(), gen_html()]
    print(f"已写入 {OUT_DIR}：")
    for p in paths:
        print(f"  {os.path.basename(p)}  ({os.path.getsize(p):,} bytes)")


if __name__ == "__main__":
    main()
